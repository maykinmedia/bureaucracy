import hashlib
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation

from bureaucracy.powerpoint import Template
from bureaucracy.powerpoint.engines import BaseEngine, PythonEngine

TEST_FILES = Path(__file__).parent / 'files'


class ConstantEngine(BaseEngine):
    def render(self, fragment, context):
        return 'Constant'


class ContextObject:
    def __init__(self, link, desc):
        self.link = link
        self.desc = desc


def test_layouts_extraction():
    test_file = str(TEST_FILES / 'template1.pptx')
    template = Template(test_file)
    assert template.layouts == ['default', 'influencers']


def test_template_render(tmpdir):
    test_file = str(TEST_FILES / 'template1.pptx')
    template = Template(test_file)
    template.render(context={}, engine=ConstantEngine())

    outfile = str(tmpdir.join('constant-engine.pptx'))
    template.save_to(outfile)

    # check that the contents are correctly templated out
    pres = Presentation(outfile)
    assert len(pres.slides) == 1
    slide = pres.slides[0]
    assert slide.slide_layout.name == 'influencers'

    assert len(slide.placeholders) == 31 + 6
    for ph in slide.placeholders:
        assert ph.text == 'Constant'


def test_template_filled_in_placeholders(tmpdir):
    """
    Assert that filled in placeholders in the actual slide are not templated away.
    """
    test_file = str(TEST_FILES / 'empty-and-filled-in-placeholders.pptx')
    template = Template(test_file)
    context = {
        'language': 'Python',
    }
    template.render(context, engine=PythonEngine())
    outfile = str(tmpdir.join('placeholders.pptx'))
    template.save_to(outfile)

    # check that the contents are correctly templated out
    pres = Presentation(outfile)
    assert len(pres.slides) == 2

    # first slide has just plain template code in placeholders
    first_slide = pres.slides[0]
    assert len(first_slide.placeholders) == 4  # title, brand_logo, two placeholders

    ph1 = first_slide.slide_layout.placeholders[2].placeholder_format.idx
    ph2 = first_slide.slide_layout.placeholders[3].placeholder_format.idx

    # ids are fixed within a slide
    assert ph1 == 11
    assert ph2 == 12

    # compare rendered output
    assert first_slide.placeholders[ph1].text == 'A simple Python string format template'
    assert first_slide.placeholders[ph2].text == 'Another simple Python string format template'

    second_slide = pres.slides[1]
    assert first_slide.slide_layout == second_slide.slide_layout

    # compare rendered output
    assert second_slide.placeholders[ph1].text == 'Filled in placeholder – should not be replaced'
    assert second_slide.placeholders[ph2].text == 'Another simple Python string format template'


def test_control_placeholder(tmpdir):
    """
    Assert that control placeholders with zero-height are removed.

    A control placeholder can modify the context datastructure in place in
    prepration for the next placeholder(s). If it's deliberately set up to be
    zero height AND the output result is empty, it should be removed from the
    presentation alltogether to reduced clutter and confusion.
    """
    test_file = str(TEST_FILES / 'control-placeholder.pptx')
    template = Template(test_file)
    assert len(template._presentation.slides[0].placeholders) == 2

    context = {'control_placeholder_no_output': ''}
    template.render(context, engine=PythonEngine())
    outfile = str(tmpdir.join('control-ph.pptx'))
    template.save_to(outfile)

    # check that the contents are correctly templated out
    pres = Presentation(outfile)
    assert len(pres.slides) == 1

    slide = pres.slides[0]
    assert len(slide.placeholders) == 1

    ph_texts = [ph.text for ph in slide.placeholders]
    assert ph_texts == ['Click to edit Master title style']


def test_ph_ordering(tmpdir):
    """
    Assert that the placeholder template fragments are passed to the template
    engine in the right order.
    """
    test_file = str(TEST_FILES / 'ordering-placeholder.pptx')
    template = Template(test_file)
    assert len(template._presentation.slides) == 1
    context = {
        'first': 'First',
        'second': 'Second',
        'third': 'Third',
        'fourth': 'Fourth',
    }

    with patch.object(PythonEngine, 'render', return_value='some-string') as mocked_render:
        template.render(context, engine=PythonEngine())

    outfile = str(tmpdir.join('order-ph.pptx'))
    template.save_to(outfile)

    # check that the contents are correctly templated out
    pres = Presentation(outfile)
    assert len(pres.slides) == 1

    expected_calls = [
        (('{second}', context),),
        (('{third}', context),),
        (('{fourth}', context),),
        (('{first}', context),),
    ]
    assert mocked_render.call_args_list == expected_calls


def test_img_placeholder(tmpdir):
    test_file = str(TEST_FILES / 'simple_img.pptx')
    template = Template(test_file)
    assert len(template._presentation.slides[0].placeholders) == 1

    goat = str(TEST_FILES / 'goat.jpg')

    context = {
        'goat_here_pls': goat,
    }

    template.render(context, engine=PythonEngine())
    outfile = str(tmpdir.join('placeholders.pptx'))
    template.save_to(outfile)

    # check that the contents are correctly templated out

    pres = Presentation(outfile)
    assert len(pres.slides) == 1

    slide = pres.slides[0]
    # although the image file itself is placed elsewhere in the pptx zip, the image tag inside the slide tag
    # holds the image's sha1. checking that is enough for now.
    with open(goat, 'rb') as goat_file:
        expected_img_hexdigest = hashlib.sha1(goat_file.read()).hexdigest()
    assert expected_img_hexdigest == slide.shapes[0].image.sha1
