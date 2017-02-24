#!/usr/bin/env python
import random
from collections import OrderedDict

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


filepath = '/home/bbt/Downloads/imtool export test 1.pptx'
save_filepath = '/home/bbt/Downloads/imtool export generated.pptx'

goat = '/home/bbt/Downloads/goat.jpg'


def construct_template_code(slide):
    """
    Build the entire template code to be rendered as one string.

    Used for debugging purposes now, to investigate the structure of
    placeholders.
    """
    template_bits = OrderedDict()

    for placeholder in slide.slide_layout.placeholders:
        template_bits[placeholder.placeholder_format.idx] = placeholder.text

    for placeholder in slide.placeholders:
        # check if there is a text already. If there is, it's an override and
        # the template bit should not be used
        if not placeholder.text:
            continue
        template_bits[placeholder.placeholder_format.idx] = placeholder.text

    print("\n".join(template_bits.values()))

    return template_bits


class SimpleEngine(object):

    words = 'I am a really, really simple engine'.split()

    def render(self, context):
        return random.choice(self.words)


class TemplateInterface(object):

    def __init__(self, engine, context):
        self.engine = engine
        self.context = context

    def render(self, fragment):
        """
        Render a template fragment with the configured template engine.
        """
        return self.engine.render(self.context)


def main():
    pres = Presentation(filepath)

    context = {
        'brand_logo': goat,
        'cards': [{
            'nickname': 'Card 1',
            'country_code': 'EST',
            'n_paid': 2,
            'paid_posts': [{'id': 1}, {'id': 2}],
            'n_earned': 0,
            'earned_posts': [],
        }]
    }
    interface = TemplateInterface(SimpleEngine(), context)

    print("Checking layouts...")
    for layout in pres.slide_layouts:
        print("  * {0.name}".format(layout))

    slides = pres.slides
    print("Presentation has {0} slide(s)".format(len(slides)))

    print("Templating out first slide...")
    slide0 = slides[0]

    print("Found placeholders:")
    for placeholder in slide0.placeholders:
        print("  * {0.placeholder_format.idx} - {0.name}".format(placeholder))

    print("Setting a test title...")
    slide0.shapes.title.text = "Example title filled in"

    template_bits = construct_template_code(slide0)
    for idx, template_bit in template_bits.items():
        placeholder = slide0.placeholders[idx]
        rendered = interface.render(template_bit)
        placeholder.text = rendered

    # print("Filling in an image")
    # logo_ph = slide0.placeholders[10]
    # assert logo_ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE, 'Not a picture placeholder'

    pres.save(save_filepath)


if __name__ == '__main__':
    main()
