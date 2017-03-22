"""
Public interface to use powerpoint presentations as export template.
"""
from io import BytesIO
from pptx import Presentation

from .engines import PythonEngine
from .slides import SlideContainer

__all__ = ['Template']


class TemplateIterator:
    """
    Iterater that knows how to deal with newly inserted slides.

    This is private API.
    """

    def __init__(self, slides):
        self.slides = slides
        self.current = 0

    def __next__(self):
        """
        Return the next slide in the slideset, which may have been inserted.
        """
        if self.current >= len(self.slides):
            raise StopIteration
        slide = self.slides[self.current]
        self.current += 1
        return slide


class Template:
    """
    A powerpoint presentation that serves as a template.

    :param filepath: path to the powerpoint file on disk or filelike object.
    """

    def __init__(self, pptx):
        self._presentation = Presentation(pptx)

    def __iter__(self):
        return TemplateIterator(self._presentation.slides)

    @property
    def layouts(self):
        """
        Returns the names of the slide layouts present in the template file.
        """
        return [layout.name for layout in self._presentation.slide_layouts]

    def render(self, context, engine=PythonEngine()):

        for slide in self:
            slide = SlideContainer(slide, self._presentation)
            slide.render(engine, context)

    def save_to(self, outfile):
        self._presentation.save(outfile)

    def to_bytes(self):
        handle = BytesIO()
        self._presentation.save(handle)
        handle.seek(0)
        return handle.read()
