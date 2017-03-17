from collections import OrderedDict

from pptx import Presentation
from pptx.slide import Slide

from bureaucracy.powerpoint.exceptions import TemplateSyntaxError

from .engines import BaseEngine
from .placeholders import PlaceholderContainer
from .shapes import ShapeContainer

CONTEXT_KEY_FOR_SLIDE = 'PPT_CURRENT_SLIDE'


class SlideContainer:
    def __init__(self, slide: Slide, presentation: Presentation):
        self.slide = slide
        self.presentation = presentation

    def __getattr__(self, name):
        """
        Proxy unknown attributes to the underlying slide object
        """
        return getattr(self.slide, name)

    def extract_shapes(self):
        """
        Extracts all the shapes from the slide and build a tree.

        Big shapes can wrap small shapes, and based on this we can re-group
        and nest them to apply an ordering to shapes (and thus placeholders).
        """
        if not self.slide.shapes:
            return []

        shapes = sorted(self.slide.slide_layout.shapes, key=lambda s: (s.width, s.height), reverse=True)
        wrapped_shapes = [ShapeContainer(shape) for shape in shapes]

        all_shapes = []

        while wrapped_shapes:
            shape = wrapped_shapes.pop(0)
            all_shapes.append(shape)
            for other_shape in wrapped_shapes:
                if not shape.wraps(other_shape):
                    continue
                shape.add_child(other_shape)

        # root shapes are the shapes without parent (= the biggest shapes that wrap other shapes)
        root_shapes = [shape for shape in all_shapes if shape.is_root]
        # finally, order by the center point of the shape
        root_shapes = sorted(root_shapes, key=lambda s: (s.center_y, s.center_x))
        return root_shapes

    def get_placeholder_idx_in_correct_order(self, fragments):
        """
        Given a slide, determine the order of template fragment evaluation.

        The slide placeholders are grouped by shapes which indicate that a
        set of placeholders needs to be evaluated before another set. This
        nesting translates into a deterministic order - top to bottom, and
        within a horizontal row from left to right.
        """
        shapes = self.extract_shapes()
        # we now have the correct order for the placeholders
        placeholders = sum((shape.get_placeholders() for shape in shapes), [])
        ordered_phs = [ph.placeholder_format.idx for ph in placeholders]
        return [idx for idx in ordered_phs if idx in fragments]

    def extract_template_code(self):
        """
        Extract the template code from slide placeholders.

        Placeholders have multiple 'levels': a placeholder can exist on a
        slide layout and on the slide itself. If the placeholder is filled in
        on the slide itself, it is not considered to be template code. If the
        value on the slide itself is empty, the value from the layout is taken
        and assumed to be template code.

        :return: an OrderedDict with placeholder id's as key and template code
          as value.
        """
        fragments = {}

        # set up the slide layout placeholder as template code
        for placeholder in self.slide.slide_layout.placeholders:
            fragments[placeholder.placeholder_format.idx] = placeholder.text

        # if a value exists for the placeholder in the slide itself, ignore the
        # template code
        for placeholder in self.slide.placeholders:
            if not placeholder.text:
                continue
            del fragments[placeholder.placeholder_format.idx]

        idxes = self.get_placeholder_idx_in_correct_order(fragments)
        # return the template bits in the right order
        return OrderedDict((idx, fragments[idx]) for idx in idxes)

    def render(self, engine: BaseEngine, context: dict):
        """
        Delegate the rendering to the underlying placeholders.
        """
        context[CONTEXT_KEY_FOR_SLIDE] = self

        fragments = self.extract_template_code()
        for idx, fragment in fragments.items():
            try:
                placeholder = PlaceholderContainer(self.slide.placeholders[idx], fragment)
            except KeyError:
                raise TemplateSyntaxError('Altough it is present on the slide master, this placeholder does'
                                          'not seem to appear on the slide iself. Did you forget to apply the layout?')
            placeholder.render(engine, context)
            if placeholder.is_empty:
                placeholder.remove()

    def insert_another(self):
        """
        Inserts the same slide into the presentation after the current position.

        NOTE: there's no insert slide method, only append to the end of the
        presentation, so we're using private API here.
        """
        layout = self.slide.slide_layout
        self.presentation.slides.add_slide(layout)

        current_index = self.presentation.slides.index(self.slide)
        # always at the end, and the ID is fixed (does not indicate position)
        # move the element to behind the current slide
        slide_list = self.presentation.slides._sldIdLst
        slide_list[current_index + 1] = slide_list[-1]
